"""Extract text, titles, tables, chart titles, and notes from a PPTX deck."""

from pptx import Presentation


def parse_pptx(path):
    """Parse a PPTX deck.

    Returns a list of dicts, one per slide:
        [{"slide": 1, "title": "...", "text": "...", "notes": "...", "images": []}]

    Note: python-pptx cannot render slide thumbnails, so image extraction of the
    rendered slide is skipped gracefully (images list is always empty). Embedded
    picture shapes could be exported, but they are not needed for text-based AI
    extraction, so we intentionally skip them to keep the parser dependency-light.
    """
    try:
        prs = Presentation(path)
    except Exception as exc:  # noqa: BLE001 - clean message for the CLI
        raise RuntimeError(
            f"Failed to open PPTX '{path}': {exc}. "
            "Note: legacy binary .ppt files are not supported."
        ) from exc

    slides = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        title = _extract_title(slide)
        text_chunks = _extract_text(slide)
        notes = _extract_notes(slide)

        slides.append(
            {
                "slide": slide_index,
                "title": title,
                "text": "\n".join(text_chunks).strip(),
                "notes": notes,
                "images": [],
            }
        )

    return slides


def _extract_title(slide):
    """Return the slide's title text if a title placeholder exists."""
    try:
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            return slide.shapes.title.text.strip()
    except Exception:  # noqa: BLE001
        pass
    return ""


def _extract_text(slide):
    """Collect text from every shape: text frames, tables, and chart titles."""
    chunks = []
    for shape in slide.shapes:
        # Text frames (title, body, text boxes)
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                chunks.append(txt)

        # Tables: flatten each row into a tab-joined line
        if shape.has_table:
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                line = "\t".join(c for c in cells if c)
                if line:
                    chunks.append(line)

        # Charts: grab the chart title if present
        if getattr(shape, "has_chart", False):
            try:
                chart = shape.chart
                if chart.has_title and chart.chart_title.has_text_frame:
                    ctitle = chart.chart_title.text_frame.text.strip()
                    if ctitle:
                        chunks.append(f"[Chart] {ctitle}")
            except Exception:  # noqa: BLE001
                pass

    return chunks


def _extract_notes(slide):
    """Return the speaker-notes text for a slide, if any."""
    try:
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                return notes_frame.text.strip()
    except Exception:  # noqa: BLE001
        pass
    return ""
